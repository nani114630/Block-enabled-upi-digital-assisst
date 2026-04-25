from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

DARK_BG = RGBColor(15, 23, 42)
PRIMARY = RGBColor(56, 189, 248)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(226, 232, 240)
DARK_GRAY = RGBColor(30, 41, 59)

def apply_background(slide, color=DARK_BG):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_header_footer(slide):
    header = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(12.333), Inches(0.4))
    tf = header.text_frame
    p = tf.paragraphs[0]
    p.text = "Nani's Project - Review 2"
    p.font.size = Pt(12)
    p.font.color.rgb = PRIMARY
    p.alignment = PP_ALIGN.RIGHT

    footer = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(12.333), Inches(0.4))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = "Nani's Project - Review 2"
    p.font.size = Pt(12)
    p.font.color.rgb = PRIMARY
    p.alignment = PP_ALIGN.CENTER

def create_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide, DARK_BG)
    
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.5), Inches(13.333), Inches(0.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY
    shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.333), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.bold = True
    p.font.size = Pt(54)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11.333), Inches(1.5))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle_text
    p.font.size = Pt(28)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER
    add_header_footer(slide)

def create_content_slide(prs, title_text, bullet_points):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide, DARK_BG)
    
    header_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = DARK_GRAY
    header_shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.bold = True
    p.font.size = Pt(36)
    p.font.color.rgb = PRIMARY
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
    tf = content_box.text_frame
    
    for pt in bullet_points:
        p = tf.add_paragraph()
        p.text = pt
        p.font.size = Pt(24)
        p.font.color.rgb = WHITE
        p.space_before = Pt(14)
    add_header_footer(slide)

def create_image_slide(prs, title_text, image_path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide, DARK_BG)
    
    header_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = DARK_GRAY
    header_shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.bold = True
    p.font.size = Pt(36)
    p.font.color.rgb = PRIMARY
    
    try:
        slide.shapes.add_picture(image_path, Inches(1), Inches(1.5), width=Inches(11.333))
    except Exception:
        pass
    add_header_footer(slide)

# 1
create_title_slide(prs, "PROJECT REVIEW 2", "UPI Digital Asset System:\nBlockchain-based Digital Asset Transparency Platform")

# 2
create_content_slide(prs, "1. Introduction", [
    "• Digital assets (NFTs) prove immutable ownership of unique items.",
    "• Current adoption is hindered by the complexity of crypto wallets.",
    "• This project proposes a hybrid marketplace bridging fiat and blockchain.",
    "• Users can purchase NFTs using traditional payment methods (UPI, Cards).",
    "• Ownership is minted and verified on the Polygon blockchain."
])

# 3
create_content_slide(prs, "2. Problem Statement", [
    "• Current NFT platforms require users to understand crypto wallets.",
    "• Users must purchase Ethereum or MATIC from exchanges first.",
    "• Managing seed phrases introduces massive security risks for beginners.",
    "• There is a lack of regulatory-friendly fiat ramps directly into NFTs."
])

# 4
create_content_slide(prs, "3. Existing System", [
    "• Conventional NFT platforms require MetaMask or similar wallets.",
    "• Reliance on Layer-1 networks like Ethereum results in high gas fees.",
    "• Fiat on-ramps involve high third-party processing fees.",
    "• Storage solutions often use centralized servers, risking broken links."
])

# 5
create_content_slide(prs, "4. Disadvantages of Existing System", [
    "• High barrier to entry for non-technical users.",
    "• Exposure to cryptocurrency price volatility.",
    "• Complex onboarding process (KYC on exchanges, wallet setup).",
    "• Increased risk of losing funds due to wallet mismanagement.",
    "• Transactions are entirely crypto-dependent."
])

# 6
create_content_slide(prs, "5. Proposed System", [
    "• Implement a full-stack Next.js and Node.js marketplace.",
    "• Integrate Razorpay for seamless INR (fiat) transactions.",
    "• Develop ERC-721 Smart Contracts on the Polygon network.",
    "• Automatically trigger NFT minting to the user's custodial address.",
    "• Store metadata on IPFS (Pinata) and off-chain data in MongoDB."
])

# 7
create_content_slide(prs, "6. Advantages of Proposed System", [
    "• Zero Crypto Knowledge Required: Users pay with familiar UPI/Cards.",
    "• Immutable Proof of Ownership: Assets are securely backed by Polygon.",
    "• Low Transaction Costs: Leveraging Polygon's Layer-2 scalability.",
    "• Fast Confirmations: Immediate fiat verification and near-instant minting.",
    "• Seamless User Experience: Traditional e-commerce feel."
])

# 8
create_content_slide(prs, "7. Project Objectives", [
    "• Simplify the UX of purchasing digital assets.",
    "• Provide blockchain transparency without the technical overhead.",
    "• Ensure security and trust using PCI-DSS compliant gateways.",
    "• Build a scalable, production-ready Web3 application."
])

# 9
create_content_slide(prs, "8. Feasibility Study: Technical", [
    "• The required tech stack (Next.js, Node.js, Polygon) is mature.",
    "• Razorpay provides robust documentation for fiat integration.",
    "• Developer tools like Hardhat and Ethers.js streamline blockchain work.",
    "• Conclusion: Technically highly feasible."
])

# 10
create_content_slide(prs, "9. Feasibility Study: Economic", [
    "• Uses open-source technologies, minimizing software costs.",
    "• Polygon Mumbai testnet allows free smart contract deployment.",
    "• IPFS Pinata offers a free tier for metadata storage.",
    "• Conclusion: Economically feasible and cost-effective."
])

# 11
create_content_slide(prs, "10. Feasibility Study: Operational", [
    "• The UI is designed to mimic standard e-commerce platforms.",
    "• Requires zero new training for the end-user.",
    "• Automated minting reduces operational overhead for administrators.",
    "• Conclusion: Operationally sound and user-friendly."
])

# 12 - BIG ARCHITECTURE IMAGE
create_image_slide(prs, "11. System Architecture Diagram", "architecture.png")

# 13
create_content_slide(prs, "12. Hardware Requirements", [
    "• Processor: Intel Core i3 or equivalent (minimum).",
    "• RAM: 8 GB or higher (recommended for development).",
    "• Storage: 256 GB SSD.",
    "• Internet: Stable broadband connection for blockchain interaction."
])

# 14
create_content_slide(prs, "13. Software Requirements", [
    "• Operating System: Windows 10/11, macOS, or Linux.",
    "• Runtime Environment: Node.js (v18+).",
    "• Database: MongoDB Atlas (Cloud).",
    "• Code Editor: Visual Studio Code.",
    "• Browser: Google Chrome or Mozilla Firefox."
])

# 15
create_content_slide(prs, "14. Technologies Used: Frontend", [
    "• Next.js 14: React framework for hybrid static & server rendering.",
    "• Tailwind CSS: Utility-first CSS framework for rapid styling.",
    "• Zustand: Lightweight global state management.",
    "• Axios: Promise-based HTTP client for the browser."
])

# 16
create_content_slide(prs, "15. Technologies Used: Backend", [
    "• Node.js: Asynchronous, event-driven JavaScript runtime.",
    "• Express.js: Fast, unopinionated web framework for Node.",
    "• JWT (JSON Web Tokens): For secure, stateless authentication.",
    "• Razorpay SDK: For generating and verifying fiat payment orders."
])

# 17
create_content_slide(prs, "16. Technologies Used: Database", [
    "• MongoDB: NoSQL database for flexible schema design.",
    "• Mongoose: Elegant MongoDB object modeling for Node.js.",
    "• Stores user profiles, asset metadata, and transaction history.",
    "• Syncs off-chain data with on-chain events."
])

# 18
create_content_slide(prs, "17. Technologies Used: Blockchain", [
    "• Solidity (0.8.20): Object-oriented programming language for smart contracts.",
    "• Hardhat: Ethereum development environment for compiling/testing.",
    "• Ethers.js: Library for interacting with the Ethereum/Polygon Blockchain.",
    "• Polygon Mumbai: Layer-2 test network for low-cost deployment."
])

# 19
create_content_slide(prs, "18. Technologies Used: Storage", [
    "• IPFS (InterPlanetary File System): Peer-to-peer hypermedia protocol.",
    "• Pinata: IPFS pinning service to ensure high availability.",
    "• Ensures metadata (images, JSON) is decentralized and immutable.",
    "• Prevents the 'broken link' problem common in NFTs."
])

# 20
create_content_slide(prs, "19. Modules Overview", [
    "The system is divided into five core modules:",
    "1. User Authentication & Profile Management",
    "2. Fiat Payment Gateway Integration",
    "3. Smart Contract & Automated Minting",
    "4. Decentralized Storage Integration",
    "5. Asset & Admin Management"
])

# 21
create_content_slide(prs, "20. Module 1: User Authentication", [
    "• Registration and Login functionalities.",
    "• Password encryption using bcrypt.",
    "• JWT generation for maintaining session states.",
    "• Role-based access control (Admin vs User)."
])

# 22
create_content_slide(prs, "21. Module 2: Fiat Payment Gateway", [
    "• Dynamic generation of INR prices.",
    "• Integration with Razorpay Orders API.",
    "• Secure processing of Credit Cards, Debit Cards, and UPI.",
    "• Webhook listening for 'payment.captured' events."
])

# 23
create_content_slide(prs, "22. Module 3: Smart Contract & Minting", [
    "• Deployment of an ERC-721 token contract.",
    "• Backend relayer wallet configuration.",
    "• Automated execution of the mint function post-payment.",
    "• Gas fee estimation and transaction signing via Ethers.js."
])

# 24
create_content_slide(prs, "23. Module 4: Decentralized Storage", [
    "• Integration with Pinata API.",
    "• Uploading raw image files to IPFS.",
    "• Generating and uploading corresponding JSON metadata.",
    "• Retrieving IPFS URIs to store permanently on the blockchain."
])

# 25
create_content_slide(prs, "24. Module 5: Admin Management", [
    "• Protected admin routes in the frontend.",
    "• Capabilities to upload new digital assets to the marketplace.",
    "• Dashboard to monitor sales and minting statuses.",
    "• Overall platform health monitoring."
])

# 26
create_content_slide(prs, "25. UML Diagrams", [
    "The following slides represent the architectural flow and design",
    "of the system using standard UML notations:",
    "• Use Case Diagram",
    "• Sequence Diagram",
    "• Activity Diagram",
    "• Data Flow Diagram",
    "• Class Diagram"
])

# 27
create_content_slide(prs, "26. Use Case Diagram", [
    "[ Please insert Use Case Diagram Image Here ]",
    " ",
    "Actors:",
    "1. User: Browse, Buy, View NFTs",
    "2. Admin: Upload Assets, Manage Orders",
    "3. System: Verify Payments, Mint NFTs"
])

# 28
create_content_slide(prs, "27. Sequence Diagram", [
    "[ Please insert Sequence Diagram Image Here ]",
    " ",
    "Flow:",
    "User -> Frontend -> Backend (Order) -> Razorpay (Payment) -> Backend (Verify) -> Polygon (Mint) -> User"
])

# 29
create_content_slide(prs, "28. Activity Diagram", [
    "[ Please insert Activity Diagram Image Here ]",
    " ",
    "Steps:",
    "Start -> Login -> Select Asset -> Initiate Payment -> Successful? -> (Yes) Mint NFT -> Update DB -> End",
    "-> (No) Show Error -> End"
])

# 30
create_content_slide(prs, "29. Data Flow Diagram (DFD)", [
    "[ Please insert DFD Image Here ]",
    " ",
    "Level 0: User interacts with UPI Digital Asset System.",
    "Level 1: Processes include Auth, Order Management, Blockchain Sync."
])

# 31
create_content_slide(prs, "30. Class Diagram", [
    "[ Please insert Class Diagram Image Here ]",
    " ",
    "Classes:",
    "• User (id, name, email, role)",
    "• Asset (id, title, price, ipfsHash)",
    "• Order (id, userId, assetId, status)",
    "• NFT (tokenId, owner, uri)"
])

# 32
create_content_slide(prs, "31. Database Design: ER Diagram", [
    "[ Please insert Entity Relationship Diagram Here ]",
    " ",
    "Relationships:",
    "• User (1) to (N) Orders",
    "• Asset (1) to (N) NFTs",
    "• Order (1) to (1) NFT"
])

# 33
create_content_slide(prs, "32. Database Schema: Users", [
    "Collection: users",
    "Fields:",
    "• _id: ObjectId",
    "• name: String",
    "• email: String (Unique)",
    "• password: String (Hashed)",
    "• role: String (Enum: 'user', 'admin')",
    "• createdAt: Date"
])

# 34
create_content_slide(prs, "33. Database Schema: Assets", [
    "Collection: assets",
    "Fields:",
    "• _id: ObjectId",
    "• title: String",
    "• description: String",
    "• priceINR: Number",
    "• imageCID: String (IPFS Hash)",
    "• metadataCID: String (IPFS Hash)",
    "• supply: Number"
])

# 35
create_content_slide(prs, "34. Database Schema: Orders", [
    "Collection: orders",
    "Fields:",
    "• _id: ObjectId",
    "• userId: ObjectId (Ref: users)",
    "• assetId: ObjectId (Ref: assets)",
    "• razorpayOrderId: String",
    "• razorpayPaymentId: String",
    "• status: String ('pending', 'paid', 'minted')"
])

# 36
create_content_slide(prs, "35. Implementation Status (Review 2)", [
    "Completed Work:",
    "• System Architecture & Database Design finalized.",
    "• Next.js Frontend UI established.",
    "• Node.js API endpoints created for Auth and Assets.",
    "• MongoDB schemas modeled and connected.",
    "• Smart Contract written and compiled."
])

# 37
create_content_slide(prs, "36. Output Screens: Homepage", [
    "[ Please insert Screenshot of Homepage ]",
    " ",
    "Displays the available digital assets fetched from the MongoDB database."
])

# 38
create_content_slide(prs, "37. Output Screens: Payment", [
    "[ Please insert Screenshot of Razorpay Popup ]",
    " ",
    "Demonstrates the fiat payment gateway integration successfully launching in test mode."
])

# 39
create_content_slide(prs, "38. Work Pending for Review 3", [
    "• Finalizing Webhook security for production.",
    "• End-to-end integration testing of the minting process.",
    "• Error handling for failed blockchain transactions.",
    "• UI polish and responsive design checks.",
    "• Deployment to Vercel (Frontend) and Render (Backend)."
])

# 40
create_content_slide(prs, "39. Conclusion", [
    "• The hybrid approach successfully abstracts Web3 complexity.",
    "• Utilizing Razorpay ensures trust and familiarity for users.",
    "• The architecture is robust, scalable, and decentralized where it matters.",
    "• The project is on track for successful completion by Review 3."
])

# 41
create_content_slide(prs, "40. References", [
    "1. Polygon Documentation: https://wiki.polygon.technology/",
    "2. Next.js Documentation: https://nextjs.org/docs",
    "3. Razorpay API Docs: https://razorpay.com/docs/api/",
    "4. IPFS & Pinata: https://docs.pinata.cloud/",
    "5. Solidity ERC-721: https://docs.openzeppelin.com/contracts/4.x/erc721"
])

# 42
create_title_slide(prs, "Thank You", "Questions & Answers")

prs.save(r"C:\Users\kande\Desktop\Nanis_Project_Review_2.pptx")
print("Review 2 Presentation saved successfully!")
