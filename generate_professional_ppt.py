from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

DARK_BG = RGBColor(15, 23, 42)
PRIMARY = RGBColor(56, 189, 248)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(226, 232, 240)
DARK_GRAY = RGBColor(30, 41, 59)

def apply_background(slide, color=DARK_BG):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_header_footer(slide):
    # Header
    header = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(12.333), Inches(0.4))
    tf = header.text_frame
    p = tf.paragraphs[0]
    p.text = "Nani's Project"
    p.font.size = Pt(12)
    p.font.color.rgb = PRIMARY
    p.alignment = PP_ALIGN.RIGHT

    # Footer
    footer = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(12.333), Inches(0.4))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = "Nani's Project"
    p.font.size = Pt(12)
    p.font.color.rgb = PRIMARY
    p.alignment = PP_ALIGN.CENTER

def create_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide, DARK_BG)
    
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.5), Inches(13.333), Inches(0.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY
    shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.333), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.bold = True
    p.font.size = Pt(54)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11.333), Inches(1.5))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle_text
    p.font.size = Pt(28)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER
    
    add_header_footer(slide)

def create_content_slide(prs, title_text, bullet_points):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide, DARK_BG)
    
    header_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = DARK_GRAY
    header_shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.bold = True
    p.font.size = Pt(36)
    p.font.color.rgb = PRIMARY
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
    tf = content_box.text_frame
    
    for pt in bullet_points:
        p = tf.add_paragraph()
        p.text = pt
        p.font.size = Pt(24)
        p.font.color.rgb = WHITE
        p.space_before = Pt(14)
        
    add_header_footer(slide)

def create_image_slide(prs, title_text, image_path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide, DARK_BG)
    
    header_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = DARK_GRAY
    header_shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.bold = True
    p.font.size = Pt(36)
    p.font.color.rgb = PRIMARY
    
    try:
        # Fill the remaining space with the image
        slide.shapes.add_picture(image_path, Inches(1), Inches(1.5), width=Inches(11.333))
    except Exception:
        pass
        
    add_header_footer(slide)

# 1
create_title_slide(prs, "UPI Digital Asset System", "Blockchain-based Digital Asset Transparency Platform\nwith INR Payment Integration")

# 2
create_content_slide(prs, "1. Abstract", [
    "• The digital asset landscape is heavily reliant on cryptocurrency wallets.",
    "• This project proposes a hybrid marketplace bridging fiat and blockchain.",
    "• Users can purchase NFTs using traditional payment methods (UPI, Cards, Net Banking).",
    "• Ownership is minted and verified on the Polygon blockchain.",
    "• Eliminates the steep learning curve of acquiring cryptocurrencies."
])

# 3
create_content_slide(prs, "2. Introduction to Web3", [
    "• Web3 represents the decentralized phase of the internet.",
    "• Relies heavily on blockchain technology for transparency and ownership.",
    "• Shifting power from centralized entities back to individual users.",
    "• Cryptocurrencies have been the primary method of interacting with Web3."
])

# 4
create_content_slide(prs, "3. Introduction to Digital Assets", [
    "• Digital assets (NFTs) prove immutable ownership of unique items.",
    "• Used for art, real estate, software licenses, and digital identity.",
    "• Stored on a decentralized ledger to prevent tampering.",
    "• Represents a massive shift in how value is stored digitally."
])

# 5
create_content_slide(prs, "4. The Problem with Current NFT Platforms", [
    "• Current adoption is severely hindered by the complexity of crypto wallets.",
    "• Users must register on exchanges, pass KYC, and buy volatile tokens.",
    "• This process alienates 99% of traditional e-commerce buyers.",
    "• A simpler bridge between Web2 (fiat) and Web3 (blockchain) is needed."
])

# 6
create_content_slide(prs, "5. Existing System: Wallet Complexities", [
    "• Conventional NFT platforms require MetaMask, Phantom, or similar wallets.",
    "• Managing seed phrases introduces massive security risks for beginners.",
    "• Gas fees are unpredictable and must be paid in native tokens.",
    "• If a user loses a seed phrase, their assets are lost forever."
])

# 7
create_content_slide(prs, "6. Existing System: Fiat On-Ramps", [
    "• Current 'fiat on-ramps' still require users to mint tokens themselves.",
    "• Often involve high third-party processing fees (e.g., MoonPay).",
    "• Still requires the user to maintain and connect a non-custodial wallet.",
    "• The user journey remains fundamentally disjointed."
])

# 8
create_content_slide(prs, "7. Disadvantages of Existing Systems (General)", [
    "• High barrier to entry for non-technical users.",
    "• Direct exposure to cryptocurrency market price volatility.",
    "• Significant drop-off rates during the checkout process.",
    "• Lack of regulatory-friendly fiat ramps for local markets like India."
])

# 9
create_content_slide(prs, "8. Disadvantages of Existing Systems (Technical)", [
    "• High transaction fees on Layer-1 networks like Ethereum.",
    "• Network congestion leads to failed minting transactions.",
    "• Decentralized storage is often not implemented correctly (broken links).",
    "• Hard for traditional businesses to integrate easily into their apps."
])

# 10
create_content_slide(prs, "9. Proposed System Overview", [
    "• A comprehensive digital marketplace that feels like standard e-commerce.",
    "• Integrates Razorpay for familiar local INR transactions.",
    "• Automates the smart contract execution securely on the backend.",
    "• Provides a transparent, scalable, and user-friendly Web3 solution."
])

# 11
create_content_slide(prs, "10. Proposed Methodology: Hybrid Marketplace", [
    "• Implement a full-stack Next.js and Node.js marketplace.",
    "• Bridge the gap by holding custodial wallets or mapping user IDs.",
    "• Display assets clearly with INR pricing, abstracted from crypto values.",
    "• Allow users to track their asset history directly on the platform."
])

# 12
create_content_slide(prs, "11. Proposed Methodology: Checkout Flow", [
    "• User clicks 'Buy' -> Backend generates a Razorpay Order.",
    "• Razorpay Popup accepts standard UPI/Card/Netbanking.",
    "• Upon payment success, a secure webhook is triggered.",
    "• The webhook securely calls the minting function on the blockchain."
])

# 13
create_content_slide(prs, "12. Proposed Methodology: Smart Contracts", [
    "• Develop ERC-721 Standard Smart Contracts.",
    "• Deploy to the Polygon network for ultra-low transaction costs.",
    "• Backend application signs the transaction using a relayer wallet.",
    "• The newly minted NFT is assigned securely to the user's mapped address."
])

# 14
create_content_slide(prs, "13. Advantages of Proposed System", [
    "• Zero Crypto Knowledge Required.",
    "• Users pay with familiar UPI/Cards.",
    "• Eliminates the need for end-users to manage gas fees.",
    "• Transactions are legally and financially compliant using standard gateways."
])

# 15
create_content_slide(prs, "14. Key Objective: Simplifying UX", [
    "• Create an interface indistinguishable from Amazon or Flipkart.",
    "• Users login via standard email/password (JWT authentication).",
    "• Dashboards clearly outline 'My Owned Assets'.",
    "• Abstract away complicated transaction hashes into clickable links."
])

# 16
create_content_slide(prs, "15. Key Objective: Blockchain Transparency", [
    "• Despite the Web2 UI, the core asset is genuinely decentralized.",
    "• Immutable Proof of Ownership backed by Polygon.",
    "• Anyone can audit the smart contract publicly on PolygonScan.",
    "• Combines the ease of centralized UI with decentralized trust."
])

# 17
create_content_slide(prs, "16. Key Objective: Security & Trust", [
    "• Payments are handled by PCI-DSS compliant Razorpay.",
    "• Smart contracts are locked and prevent unauthorized minting.",
    "• JWT tokens prevent unauthorized access to the application API.",
    "• Double verification mechanisms (MongoDB record + Blockchain record)."
])

# 18
create_content_slide(prs, "17. Tech Stack Overview", [
    "• Frontend: Next.js 14, Tailwind CSS, Zustand",
    "• Backend: Node.js, Express.js, MongoDB",
    "• Blockchain: Solidity 0.8.20, Hardhat, Ethers.js",
    "• Networks & Storage: Polygon Mumbai, IPFS (Pinata)",
    "• Infrastructure: MongoDB Atlas, Vercel, Render"
])

# 19
create_content_slide(prs, "18. Frontend Technology: Next.js", [
    "• A React framework providing hybrid static & server rendering.",
    "• Enhances SEO and performance for marketplace pages.",
    "• Next.js API routes used to securely handle frontend configurations.",
    "• Enables fast, dynamic routing for individual asset pages."
])

# 20
create_content_slide(prs, "19. Frontend Technology: Tailwind & Zustand", [
    "• Tailwind CSS allows for rapid, utility-first styling.",
    "• Ensures a highly responsive and modern design system.",
    "• Zustand is used for lightweight, fast global state management.",
    "• Handles user authentication state and shopping cart data across components."
])

# 21
create_content_slide(prs, "20. Backend Technology: Node.js", [
    "• Asynchronous, event-driven JavaScript runtime.",
    "• Highly scalable for handling concurrent payment webhook requests.",
    "• Acts as the secure middleware between the frontend and the blockchain.",
    "• Prevents frontend exposure of sensitive private keys."
])

# 22
create_content_slide(prs, "21. Backend Technology: Express & MongoDB", [
    "• Express.js handles REST API routing and middleware.",
    "• MongoDB provides a flexible NoSQL schema for assets and users.",
    "• Mongoose ODM allows strict typing and relationship mapping.",
    "• MongoDB Atlas provides reliable, cloud-hosted database scalability."
])

# 23
create_content_slide(prs, "22. Database Schema: Users Collection", [
    "• Stores user credentials (hashed passwords) and profiles.",
    "• Maps traditional Web2 user accounts to internal Web3 wallets.",
    "• Stores JWT refresh tokens and role-based access levels.",
    "• Critical for linking fiat payment histories to user accounts."
])

# 24
create_content_slide(prs, "23. Database Schema: Assets & Orders", [
    "• Assets Collection: Contains metadata, pricing, and IPFS CIDs.",
    "• Orders Collection: Tracks the Razorpay order ID and payment status.",
    "• Relationships: User -> Order -> NFT -> Asset.",
    "• Ensures data consistency between pending payments and minted assets."
])

# 25
create_content_slide(prs, "24. Payment Gateway: Razorpay", [
    "• Industry-leading payment gateway for Indian users.",
    "• Supports credit/debit cards, net banking, and UPI.",
    "• Provides secure order generation (Server to Server).",
    "• Handles payment status tracking and error handling gracefully."
])

# 26
create_content_slide(prs, "25. Razorpay Webhook Handling", [
    "• The backend exposes a secure POST endpoint for Razorpay webhooks.",
    "• Webhooks verify the cryptographic signature sent by Razorpay.",
    "• Ensures users cannot spoof 'payment success' messages.",
    "• Webhook triggers the crucial `verify-mint` smart contract function."
])

# 27
create_content_slide(prs, "26. Blockchain Network: Polygon Mumbai", [
    "• A test network replicating the main Polygon PoS chain.",
    "• Used for safe, cost-free development and testing.",
    "• fully EVM compatible (Ethereum Virtual Machine).",
    "• Allows usage of standard tools like Hardhat and Ethers.js."
])

# 28
create_content_slide(prs, "27. Why Polygon? (Layer 2 Benefits)", [
    "• Scalability: Can handle thousands of transactions per second.",
    "• Low Cost: Transaction fees are typically fractions of a cent.",
    "• Fast Confirmations: Near-instant block times (2 seconds).",
    "• Perfect for an e-commerce platform requiring fast checkout experiences."
])

# 29
create_content_slide(prs, "28. Smart Contract: ERC721 Standard", [
    "• The industry standard interface for Non-Fungible Tokens.",
    "• Inherits from OpenZeppelin for battle-tested security.",
    "• Enforces uniqueness; no two tokens are exactly alike.",
    "• Allows tracking of exact ownership changes throughout history."
])

# 30
create_content_slide(prs, "29. Minting Process Explained", [
    "• The Node.js backend acts as the 'Owner' of the contract.",
    "• Calls `mintAsset(userAddress, metadataURI)` using Ethers.js.",
    "• Pays the required MATIC gas fee on behalf of the user.",
    "• Returns the blockchain Transaction Hash back to the MongoDB record."
])

# 31
create_content_slide(prs, "30. Storage Solution: IPFS", [
    "• InterPlanetary File System (IPFS).",
    "• A decentralized, peer-to-peer file sharing network.",
    "• Used to store the actual digital asset (image, video) and its JSON metadata.",
    "• Prevents the 'broken link' problem common in centralized servers."
])

# 32
create_content_slide(prs, "31. Why IPFS? (Pinata)", [
    "• Pinata acts as a pinning service to ensure files stay online permanently.",
    "• Provides extremely fast IPFS gateways for the Next.js frontend.",
    "• Blockchain only stores the IPFS hash (CID) to save gas costs.",
    "• Ensures the asset data remains as decentralized as the blockchain itself."
])

# 33 - BIG ARCHITECTURE IMAGE
create_image_slide(prs, "32. System Architecture Diagram", "architecture.png")

# 34
create_content_slide(prs, "33. Module 1: Authentication", [
    "• Secure Registration & Login endpoints.",
    "• Password hashing via bcrypt.",
    "• JWT tokens issued with configurable expiration times.",
    "• Role-based access control (Admin vs. Standard User)."
])

# 35
create_content_slide(prs, "34. Module 2: Asset Marketplace", [
    "• Dynamic fetching of active digital assets.",
    "• Filtering and search capabilities.",
    "• Detailed asset views showing IPFS metadata and current supply.",
    "• Admin panel for uploading and managing new listings."
])

# 36
create_content_slide(prs, "35. Module 3: Payment Processing", [
    "• Dynamic generation of INR prices.",
    "• Initialization of Razorpay sessions.",
    "• Handling payment callbacks and edge-case failures.",
    "• Mapping successful transaction IDs to user orders."
])

# 37
create_content_slide(prs, "36. Module 4: Blockchain Minting", [
    "• Ethers.js provider setup communicating with Polygon RPC.",
    "• Wallet initialization with securely stored private keys.",
    "• Contract interaction and gas estimation.",
    "• Post-mint synchronization with the main MongoDB database."
])

# 38
create_content_slide(prs, "37. Output Screens & Demonstration", [
    "• [ Placeholder: Homepage & Marketplace Dashboard ]",
    "• [ Placeholder: Razorpay Payment Popup in Action ]",
    "• [ Placeholder: User Profile - My NFTs Section ]",
    "• [ Placeholder: PolygonScan Verification Page ]",
    " ",
    "(Please insert actual system screenshots here)"
])

# 39
create_content_slide(prs, "38. Future Scope & Enhancements", [
    "• Direct integration with Mainnet Polygon for production use.",
    "• Support for secondary market trading (User to User sales).",
    "• Integration with physical asset tracking (Real-World Assets).",
    "• Mobile application development using React Native."
])

# 40
create_content_slide(prs, "39. Conclusion", [
    "• Successfully bridged the gap between Web2 payments and Web3 assets.",
    "• Abstracted away the complexities that hinder mainstream NFT adoption.",
    "• Proved that fiat-to-NFT minting is secure, fast, and scalable.",
    "• Created a highly robust foundation for digital asset e-commerce."
])

# 41
create_title_slide(prs, "Thank You", "Questions & Answers")

prs.save(r"C:\Users\kande\Desktop\Nanis_Project_40_Slides.pptx")
print("41-Slide Presentation saved successfully!")
