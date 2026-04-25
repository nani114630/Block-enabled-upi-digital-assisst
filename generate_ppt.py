import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

prs = Presentation()

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_box = slide.shapes.title
    subtitle_box = slide.placeholders[1]
    title_box.text = title
    subtitle_box.text = subtitle

def add_bullet_slide(prs, title, points):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title_box = slide.shapes.title
    title_box.text = title
    content_box = slide.placeholders[1]
    tf = content_box.text_frame
    tf.text = points[0]
    for point in points[1:]:
        p = tf.add_paragraph()
        p.text = point

# Slide 1: Title
add_title_slide(prs, "UPI Digital Asset System", "Blockchain-based Digital Asset Transparency Platform\nwith INR Payment Integration")

# Slide 2: Overview & Problem Statement
add_bullet_slide(prs, "Overview & Solution", [
    "Problem: Digital assets usually require crypto wallets and complex cryptocurrency transactions.",
    "Solution: A system bridging traditional payments with blockchain technology.",
    "Users can pay in INR using familiar methods (UPI, Card, Net Banking) via Razorpay.",
    "NFTs are minted on Polygon, providing immutable proof of ownership.",
    "Creates a seamless experience for non-crypto users while maintaining blockchain transparency."
])

# Slide 3: Core Features
add_bullet_slide(prs, "Core Features", [
    "User Authentication: Secure JWT-based auth with role management.",
    "Asset Marketplace: Browse, search, and filter digital assets.",
    "INR Payments: Real-time fiat payments via Razorpay.",
    "Automated NFT Minting: Automatic minting after successful payment.",
    "Dual Verification: Ownership verified via MongoDB and Polygon.",
    "Ownership History: Full transaction history stored in the smart contract."
])

# Slide 4: Tech Stack
add_bullet_slide(prs, "Tech Stack", [
    "Frontend: Next.js 14, Tailwind CSS, Zustand",
    "Backend: Node.js, Express.js, MongoDB, Mongoose",
    "Blockchain: Solidity 0.8.20, Hardhat, Ethers.js",
    "Networks & Storage: Polygon Mumbai Testnet, IPFS (Pinata)",
    "Infrastructure: MongoDB Atlas, Vercel, Render"
])

# Slide 5: Payment Flow
add_bullet_slide(prs, "End-to-End Payment Flow", [
    "1. USER selects an asset from the marketplace.",
    "2. BACKEND creates an INR order and generates a Razorpay session.",
    "3. USER completes the payment using traditional methods (e.g., Cards, UPI).",
    "4. BACKEND receives the webhook, verifying the payment status.",
    "5. SMART CONTRACT is triggered to mint the NFT on Polygon.",
    "6. USER immediately sees the NFT in their dashboard with a PolygonScan link."
])

# Slide 6: Blockchain Integration
add_bullet_slide(prs, "Blockchain Integration", [
    "Smart Contract (AssetNFT.sol) is an ERC721 token.",
    "Uses Polygon Network for low fees (~$0.001) and fast confirmations.",
    "EVM compatibility allows interaction with standard Ethereum tools.",
    "Complete transparency - Ownership is publicly verifiable on PolygonScan."
])

# Slide 7: Database Architecture
add_bullet_slide(prs, "Database Architecture", [
    "Collections Structure:",
    " - Users: Account details and authentication",
    " - Assets: Digital assets listed in the marketplace",
    " - NFTs: Records of minted non-fungible tokens",
    " - Orders: Purchase orders initiated by users",
    " - Transactions: History of payments and status",
    "Relationship: User -> Order -> NFT -> Asset"
])

# Slide 8: Thank You
add_title_slide(prs, "Thank You", "Questions & Answers")

prs.save("UPI_Digital_Asset_System_Presentation.pptx")
print("Presentation generated successfully!")
